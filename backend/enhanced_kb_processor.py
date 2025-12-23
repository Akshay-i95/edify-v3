"""
Enhanced KB Multi-Format Processor
Handles PDF, PPTX, DOCX, MP4, and other formats from KB module
Organizes data into Pinecone namespaces: kb-esp, kb-psp, kb-msp, kb-ssp
Enhanced with OCR, PDF repair, and multi-method extraction from auto-ingestion
"""

import io
import logging
import time
import os
import tempfile
import warnings
from typing import List, Dict, Tuple, Optional, Any
from pathlib import Path
import asyncio
from concurrent.futures import ThreadPoolExecutor, as_completed
from tqdm import tqdm
import json
from datetime import datetime, timezone

# File format processors
import PyPDF2
import pdfplumber
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
import zipfile
import xml.etree.ElementTree as ET

# PDF Validation and Repair
try:
    import pikepdf
    PIKEPDF_AVAILABLE = True
except ImportError:
    PIKEPDF_AVAILABLE = False

# OCR and Image Processing
try:
    import pytesseract
    from PIL import Image
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

# Azure and Vector DB
from azure.storage.blob import BlobServiceClient
from azure.core.exceptions import AzureError
from pinecone import Pinecone
from sentence_transformers import SentenceTransformer

# Text processing
import tiktoken
import re

class EnhancedKBProcessor:
    def __init__(self, config: Dict):
        """Initialize the enhanced KB processor for multiple file formats"""
        self.config = config
        self.logger = logging.getLogger(__name__)
        
        # Initialize tokenizer
        try:
            self.tokenizer = tiktoken.get_encoding("cl100k_base")
        except Exception as e:
            self.logger.warning(f"⚠️ Failed to load tokenizer: {str(e)}")
            self.tokenizer = None
        
        # Initialize embedding model
        try:
            self.embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
            self.logger.info("✅ Embedding model loaded")
        except Exception as e:
            self.logger.error(f"❌ Failed to load embedding model: {str(e)}")
            raise
        
        # Configure processing parameters
        self.chunk_size = int(config.get('chunk_size', 600))
        self.chunk_overlap = int(config.get('chunk_overlap', 120))
        self.min_chunk_length = int(config.get('min_chunk_length', 100))
        self.max_chunk_length = int(config.get('max_chunk_length', 1200))
        self.batch_size = int(config.get('batch_size', 15))
        self.max_workers = int(config.get('max_workers', 4))
        
        # OCR Configuration (from auto-ingestion)
        self.enable_ocr = config.get('enable_ocr', True) and OCR_AVAILABLE
        self.ocr_language = config.get('ocr_language', 'eng')
        self.ocr_dpi = int(config.get('ocr_dpi', 300))
        self.image_to_text = config.get('image_to_text', True)
        
        # PDF Validation & Repair
        self.enable_repair = config.get('enable_repair', True) and PIKEPDF_AVAILABLE
        
        # Namespace mapping
        self.namespace_mapping = {
            'kb-esp': ['playgroup', 'ik1', 'ik2', 'ik3'],  # Early Stage Program
            'kb-psp': ['grade1', 'grade2', 'grade3', 'grade4', 'grade5'],  # Primary Stage Program
            'kb-msp': ['grade6', 'grade7', 'grade8', 'grade9', 'grade10'],  # Middle Stage Program
            'kb-ssp': ['grade11', 'grade12']  # Senior Stage Program
        }
        
        # File type handlers
        self.file_handlers = {
            '.pdf': self._process_pdf,
            '.pptx': self._process_pptx,
            '.docx': self._process_docx,
            '.mp4': self._process_mp4,
            '.ppt': self._process_ppt,
            '.xlsx': self._process_xlsx,
            '.doc': self._process_doc,
            '.html': self._process_html,
            '.pptm': self._process_pptx  # Same as pptx
        }
        
        # Statistics tracking (enhanced)
        self.stats = {
            'total_files': 0,
            'processed_files': 0,
            'failed_files': 0,
            'chunks_created': 0,
            'files_by_type': {},
            'files_by_namespace': {},
            'processing_times': {},
            'start_time': time.time(),
            'ocr_extractions': 0,
            'repair_attempts': 0,
            'successful_extractions': 0,
            'current_file': '',
            'completion_percentage': 0.0
        }
        
        # Initialize Azure and Pinecone
        self._initialize_connections()
        
    def _initialize_connections(self):
        """Initialize Azure Blob Storage and Pinecone connections"""
        try:
            # Azure Blob Storage
            connection_string = os.getenv('AZURE_STORAGE_CONNECTION_STRING')
            if not connection_string:
                raise ValueError("AZURE_STORAGE_CONNECTION_STRING not found")
            
            self.blob_service_client = BlobServiceClient.from_connection_string(connection_string)
            self.container_name = 'edifydocumentcontainer'
            self.container_client = self.blob_service_client.get_container_client(self.container_name)
            
            # Pinecone
            pinecone_api_key = os.getenv('PINECONE_API_KEY')
            if not pinecone_api_key:
                raise ValueError("PINECONE_API_KEY not found")
            
            self.pc = Pinecone(api_key=pinecone_api_key)
            
            self.index_name = os.getenv('PINECONE_INDEX_NAME', 'edify-edicenter')
            self.index = self.pc.Index(self.index_name)
            
            self.logger.info("✅ Azure and Pinecone connections initialized")
            
        except Exception as e:
            self.logger.error(f"❌ Failed to initialize connections: {str(e)}")
            raise
    
    def _validate_dependencies(self):
        """Validate OCR and PDF repair dependencies"""
        if self.enable_ocr and not OCR_AVAILABLE:
            self.logger.warning("⚠️ OCR requested but dependencies not available")
            self.enable_ocr = False
        
        if self.enable_ocr and OCR_AVAILABLE:
            try:
                version = pytesseract.get_tesseract_version()
                self.logger.info(f"✅ Tesseract OCR available: {version}")
            except Exception as e:
                self.logger.warning(f"⚠️ Tesseract not properly installed: {str(e)}")
                self.enable_ocr = False
        
        if self.enable_repair and not PIKEPDF_AVAILABLE:
            self.logger.warning("⚠️ PDF repair requested but pikepdf not available")
            self.enable_repair = False
    
    def _validate_and_repair_pdf(self, pdf_stream: io.BytesIO, blob_name: str) -> Tuple[io.BytesIO, bool]:
        """Validate and attempt to repair corrupted PDF files"""
        repaired = False
        
        try:
            # Basic validation with PyPDF2
            pdf_stream.seek(0)
            reader = PyPDF2.PdfReader(pdf_stream)
            
            # Check if we can access pages
            num_pages = len(reader.pages)
            if num_pages == 0:
                raise Exception("PDF has no pages")
            
            # Try to read first page with warning suppression
            first_page = reader.pages[0]
            with warnings.catch_warnings():
                warnings.filterwarnings("ignore")
                _ = first_page.extract_text()
            
            self.logger.debug(f"✅ PDF validation passed for {blob_name} ({num_pages} pages)")
            pdf_stream.seek(0)
            return pdf_stream, repaired
            
        except Exception as e:
            self.logger.warning(f"⚠️ PDF validation failed for {blob_name}: {str(e)}")
            
            if self.enable_repair and PIKEPDF_AVAILABLE:
                try:
                    self.stats['repair_attempts'] += 1
                    self.logger.info(f"🔧 Attempting to repair {blob_name}")
                    
                    # Save to temporary file for pikepdf
                    pdf_stream.seek(0)
                    with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_file:
                        temp_file.write(pdf_stream.getvalue())
                        temp_path = temp_file.name
                    
                    try:
                        # Use pikepdf to repair
                        with pikepdf.open(temp_path, allow_overwriting_input=True) as pdf:
                            repaired_stream = io.BytesIO()
                            pdf.save(repaired_stream)
                            repaired_stream.seek(0)
                            repaired = True
                            self.logger.info(f"✅ Successfully repaired {blob_name}")
                            return repaired_stream, repaired
                    finally:
                        os.unlink(temp_path)
                        
                except Exception as repair_error:
                    self.logger.error(f"❌ Failed to repair {blob_name}: {str(repair_error)}")
            
            pdf_stream.seek(0)
            return pdf_stream, repaired
            raise
    
    def get_namespace_for_grade(self, grade: str) -> str:
        """Determine the Pinecone namespace for a given grade"""
        grade_lower = grade.lower()
        
        for namespace, grades in self.namespace_mapping.items():
            if grade_lower in grades:
                return namespace
        
        # Default fallback
        if 'common' in grade_lower:
            return 'kb-psp'  # Common files go to primary
        
        return 'kb-msp'  # Default to middle stage
    
    def get_namespace_from_blob_path(self, blob_path: str) -> str:
        """
        Extract namespace from Azure blob path
        
        Examples:
        - kb/12/grade1/file.pdf -> kb-psp
        - kb/12/common/file.pdf -> kb-psp
        - kb/12/grade11/file.pdf -> kb-ssp
        - edipedia/2025-2026/preschools/file.pdf -> edipedia-preschools
        """
        path_lower = blob_path.lower()
        
        # Handle KB module paths: kb/<school_id>/<grade>/file.ext
        if path_lower.startswith('kb/'):
            parts = blob_path.split('/')
            if len(parts) >= 3:
                grade = parts[2].lower()  # Extract grade folder
                
                # Map grade to namespace
                for namespace, grades in self.namespace_mapping.items():
                    if grade in grades:
                        self.logger.debug(f"Mapped '{blob_path}' -> '{namespace}'")
                        return namespace
                
                # Common files default to primary
                if 'common' in grade:
                    return 'kb-psp'
        
        # Handle Edipedia paths: edipedia/<year>/<category>/file.ext
        if path_lower.startswith('edipedia/'):
            if 'preschool' in path_lower:
                return 'edipedia-preschools'
            elif 'edifyho' in path_lower or 'ho' in path_lower:
                return 'edipedia-edifyho'
            elif 'k12' in path_lower or 'k-12' in path_lower:
                return 'edipedia-k12'
            else:
                return 'edipedia-k12'  # Default for edipedia
        
        # Default fallback
        self.logger.warning(f"⚠️ Could not determine namespace for '{blob_path}', using 'kb-msp'")
        return 'kb-msp'
    
    async def process_kb_directory(self, target_year: str = '12'):
        """Process all files in the specified KB directory year"""
        self.logger.info(f"🚀 Starting KB processing for year: {target_year}")
        
        try:
            # Get all files from the target directory
            files = list(self.container_client.list_blobs(name_starts_with=f'kb/{target_year}/'))
            
            if not files:
                self.logger.warning(f"⚠️ No files found in kb/{target_year}/")
                return
            
            # Organize files by grade and type
            organized_files = self._organize_files(files, target_year)
            
            # Process each namespace separately
            for namespace, file_data in organized_files.items():
                self.logger.info(f"📚 Processing namespace: {namespace} ({len(file_data)} files)")
                await self._process_namespace(namespace, file_data)
            
            # Final statistics
            self._print_final_statistics()
            
        except Exception as e:
            self.logger.error(f"❌ Error processing KB directory: {str(e)}")
            raise
    
    def _organize_files(self, files, target_year: str) -> Dict[str, List]:
        """Organize files by namespace and collect metadata"""
        organized = {
            'kb-esp': [],
            'kb-psp': [],
            'kb-msp': [],
            'kb-ssp': []
        }
        
        for blob in files:
            try:
                path_parts = blob.name.split('/')
                if len(path_parts) < 3:
                    continue
                
                year = path_parts[1]
                grade = path_parts[2]
                filename = path_parts[-1]
                
                # Skip if not target year
                if year != target_year:
                    continue
                
                # Get file extension
                file_ext = Path(filename).suffix.lower()
                
                # Skip unsupported file types
                if file_ext not in self.file_handlers:
                    self.logger.debug(f"⚠️ Unsupported file type: {filename}")
                    continue
                
                # Determine namespace
                namespace = self.get_namespace_for_grade(grade)
                
                # Add to organized structure
                file_info = {
                    'blob_name': blob.name,
                    'filename': filename,
                    'grade': grade,
                    'file_type': file_ext,
                    'size': blob.size,
                    'last_modified': blob.last_modified
                }
                
                organized[namespace].append(file_info)
                
                # Update statistics
                if file_ext not in self.stats['files_by_type']:
                    self.stats['files_by_type'][file_ext] = 0
                self.stats['files_by_type'][file_ext] += 1
                
                if namespace not in self.stats['files_by_namespace']:
                    self.stats['files_by_namespace'][namespace] = 0
                self.stats['files_by_namespace'][namespace] += 1
                
            except Exception as e:
                self.logger.warning(f"⚠️ Error organizing file {blob.name}: {str(e)}")
                continue
        
        self.stats['total_files'] = sum(len(files) for files in organized.values())
        
        self.logger.info(f"📊 File organization complete:")
        for namespace, files in organized.items():
            self.logger.info(f"   {namespace}: {len(files)} files")
        
        return organized
    
    async def _process_namespace(self, namespace: str, files: List[Dict]):
        """Process all files for a specific namespace"""
        self.logger.info(f"🔄 Processing {len(files)} files for namespace: {namespace}")
        
        # Process files in batches
        for i in range(0, len(files), self.batch_size):
            batch = files[i:i + self.batch_size]
            
            self.logger.info(f"📦 Processing batch {i//self.batch_size + 1}/{(len(files)-1)//self.batch_size + 1}")
            
            # Process batch
            batch_chunks = self._process_file_batch(batch, namespace)
            
            # Upload to Pinecone
            if batch_chunks:
                await self._upload_to_pinecone(batch_chunks, namespace)
            
            # Update progress
            self.stats['completion_percentage'] = ((i + len(batch)) / self.stats['total_files']) * 100
            
            # Live statistics
            self._print_live_statistics()
    
    def _process_file_batch(self, files: List[Dict], namespace: str) -> List[Dict]:
        """Process a batch of files with reduced memory usage"""
        all_chunks = []
        
        # Reduce concurrent workers to prevent memory exhaustion
        max_workers = min(2, self.max_workers)  # Max 2 workers
        
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            # Process files in smaller sub-batches to manage memory
            batch_size = 2  # Process 2 files at a time
            
            for i in range(0, len(files), batch_size):
                sub_batch = files[i:i + batch_size]
                
                future_to_file = {
                    executor.submit(self._process_single_file, file_info, namespace): file_info
                    for file_info in sub_batch
                }
                
                for future in as_completed(future_to_file):
                    file_info = future_to_file[future]
                    try:
                        chunks = future.result()
                        if chunks:
                            all_chunks.extend(chunks)
                            self.stats['processed_files'] += 1
                        else:
                            self.stats['failed_files'] += 1
                        
                        # Clear chunk data immediately after processing
                        del chunks
                        
                    except Exception as e:
                        self.logger.error(f"❌ Error processing {file_info['filename']}: {str(e)}")
                        self.stats['failed_files'] += 1
                
                # Force garbage collection after each sub-batch
                import gc
                gc.collect()
        
        return all_chunks
    
    def _process_single_file(self, file_info: Dict, namespace: str) -> List[Dict]:
        """Process a single file based on its type"""
        try:
            self.stats['current_file'] = file_info['filename']
            start_time = time.time()
            
            blob_name = file_info['blob_name']
            file_type = file_info['file_type']
            
            # Get the appropriate handler
            handler = self.file_handlers.get(file_type)
            if not handler:
                self.logger.warning(f"⚠️ No handler for file type: {file_type}")
                return []
            
            # Download file content
            blob_client = self.container_client.get_blob_client(blob_name)
            file_content = blob_client.download_blob().readall()
            
            # Process with appropriate handler
            text_content, metadata = handler(file_content, file_info)
            
            if not text_content or len(text_content.strip()) < 50:
                self.logger.warning(f"⚠️ Minimal content extracted from {file_info['filename']}")
                return []
            
            # Create chunks
            chunks = self._create_chunks(text_content, file_info, metadata, namespace)
            
            # Update processing time
            processing_time = time.time() - start_time
            self.stats['processing_times'][file_type] = self.stats['processing_times'].get(file_type, [])
            self.stats['processing_times'][file_type].append(processing_time)
            
            self.logger.info(f"✅ Processed {file_info['filename']}: {len(chunks)} chunks in {processing_time:.2f}s")
            
            return chunks
            
        except Exception as e:
            self.logger.error(f"❌ Error processing {file_info['filename']}: {str(e)}")
            return []
    
    def _process_pdf(self, content: bytes, file_info: Dict) -> Tuple[str, Dict]:
        """Enhanced PDF processing with multi-method extraction and OCR"""
        try:
            pdf_stream = io.BytesIO(content)
            text_content = ""
            metadata = {
                'filename': file_info['filename'],
                'pages': 0,
                'extraction_method': 'none',
                'processing_time': 0,
                'ocr_used': False,
                'images_processed': 0,
                'repaired': False
            }
            
            start_time = time.time()
            
            # Validate and potentially repair the PDF
            pdf_stream, was_repaired = self._validate_and_repair_pdf(pdf_stream, file_info['filename'])
            if was_repaired:
                metadata['repaired'] = True
            
            # Method 1: Try pdfplumber first (best for structured text)
            try:
                pdf_stream.seek(0)
                # Suppress pdfminer warnings
                pdfminer_logger = logging.getLogger('pdfminer')
                original_level = pdfminer_logger.level
                pdfminer_logger.setLevel(logging.ERROR)
                
                try:
                    with warnings.catch_warnings():
                        warnings.filterwarnings("ignore")
                        with pdfplumber.open(pdf_stream) as pdf:
                            pages_text = []
                            for page in pdf.pages:
                                try:
                                    page_text = page.extract_text()
                                    if page_text and page_text.strip():
                                        pages_text.append(page_text)
                                except Exception as e:
                                    self.logger.debug(f"Page extraction error: {str(e)}")
                            
                            if pages_text:
                                text_content = "\n\n".join(pages_text)
                                metadata['pages'] = len(pdf.pages)
                                metadata['extraction_method'] = 'pdfplumber'
                finally:
                    pdfminer_logger.setLevel(original_level)
                
                if text_content and len(text_content.strip()) >= 100:
                    metadata['processing_time'] = time.time() - start_time
                    return self._clean_extracted_text(text_content), metadata
                    
            except Exception as e:
                self.logger.debug(f"pdfplumber failed for {file_info['filename']}: {str(e)}")
            
            # Method 2: PyMuPDF (better for complex layouts)
            if not text_content or len(text_content.strip()) < 100:
                try:
                    pdf_stream.seek(0)
                    pdf_doc = fitz.open(stream=pdf_stream.getvalue(), filetype="pdf")
                    pages_text = []
                    
                    for page_num in range(pdf_doc.page_count):
                        try:
                            page = pdf_doc[page_num]
                            page_text = page.get_text()
                            if page_text and page_text.strip():
                                pages_text.append(page_text)
                        except Exception as e:
                            self.logger.warning(f"Error extracting page {page_num}: {str(e)}")
                    
                    pdf_doc.close()
                    
                    if pages_text:
                        text_content = "\n\n".join(pages_text)
                        metadata['pages'] = len(pages_text)
                        metadata['extraction_method'] = 'PyMuPDF'
                        
                except Exception as e:
                    self.logger.warning(f"PyMuPDF failed for {file_info['filename']}: {str(e)}")
            
            # Method 3: OCR for scanned documents
            if (len(text_content.strip()) < 100 and self.enable_ocr):
                self.logger.info(f"🔍 Trying OCR for {file_info['filename']}")
                ocr_text, ocr_metadata = self._extract_with_ocr(pdf_stream, file_info['filename'])
                
                if len(ocr_text.strip()) > len(text_content.strip()):
                    text_content = ocr_text
                    metadata.update(ocr_metadata)
                    metadata['ocr_used'] = True
                    self.stats['ocr_extractions'] += 1
            
            # Extract images and convert to text if enabled
            if self.image_to_text and self.enable_ocr and len(text_content.strip()) >= 100:
                image_text, image_count = self._extract_images_to_text(pdf_stream, file_info['filename'])
                if image_text.strip():
                    text_content += "\n\n--- IMAGE CONTENT ---\n" + image_text
                    metadata['images_processed'] = image_count
            
            # Clean text
            if text_content:
                text_content = self._clean_extracted_text(text_content)
            
            metadata['processing_time'] = time.time() - start_time
            metadata['character_count'] = len(text_content)
            
            self.logger.info(f"Extracted {len(text_content)} chars from {file_info['filename']} using {metadata['extraction_method']}")
            
            return text_content, metadata
            
        except Exception as e:
            self.logger.error(f"❌ PDF processing failed: {str(e)}")
            return "", {}
    
    def _extract_with_ocr(self, pdf_stream: io.BytesIO, filename: str) -> Tuple[str, Dict]:
        """Extract text using OCR for scanned documents"""
        if not self.enable_ocr or not OCR_AVAILABLE:
            return "", {}
        
        text_content = ""
        metadata = {'extraction_method': 'OCR', 'pages': 0}
        
        try:
            pdf_stream.seek(0)
            pdf_doc = fitz.open(stream=pdf_stream.getvalue(), filetype="pdf")
            pages_text = []
            
            self.logger.info(f"🔍 Starting OCR processing for {filename} ({pdf_doc.page_count} pages)")
            
            for page_num in range(min(pdf_doc.page_count, 50)):  # Limit OCR to 50 pages
                try:
                    page = pdf_doc[page_num]
                    
                    # Convert page to image
                    pix = page.get_pixmap(matrix=fitz.Matrix(self.ocr_dpi/72, self.ocr_dpi/72))
                    img_data = pix.tobytes("png")
                    image = Image.open(io.BytesIO(img_data))
                    
                    # Apply OCR
                    page_text = pytesseract.image_to_string(
                        image, 
                        lang=self.ocr_language,
                        config='--psm 6'
                    )
                    
                    if page_text and page_text.strip():
                        pages_text.append(page_text)
                    
                    # Memory cleanup
                    image.close()
                    pix = None
                    
                except Exception as e:
                    self.logger.warning(f"⚠️ OCR failed for page {page_num}: {str(e)}")
                    continue
            
            pdf_doc.close()
            
            if pages_text:
                text_content = "\n\n".join(pages_text)
                metadata['pages'] = len(pages_text)
                self.logger.info(f"✅ OCR extracted {len(text_content)} chars from {filename}")
            
        except Exception as e:
            self.logger.error(f"❌ OCR extraction failed: {str(e)}")
        
        return text_content, metadata
    
    def _extract_images_to_text(self, pdf_stream: io.BytesIO, filename: str) -> Tuple[str, int]:
        """Extract text from images within PDF"""
        if not self.enable_ocr or not OCR_AVAILABLE:
            return "", 0
        
        image_texts = []
        image_count = 0
        
        try:
            pdf_stream.seek(0)
            pdf_doc = fitz.open(stream=pdf_stream.getvalue(), filetype="pdf")
            
            for page_num in range(pdf_doc.page_count):
                try:
                    page = pdf_doc[page_num]
                    image_list = page.get_images()
                    
                    for img_index, img in enumerate(image_list):
                        try:
                            xref = img[0]
                            base_image = pdf_doc.extract_image(xref)
                            image_bytes = base_image["image"]
                            
                            # Convert to PIL Image
                            image = Image.open(io.BytesIO(image_bytes))
                            
                            # Apply OCR
                            img_text = pytesseract.image_to_string(image, lang=self.ocr_language)
                            if img_text and len(img_text.strip()) > 20:
                                image_texts.append(img_text.strip())
                                image_count += 1
                            
                            image.close()
                            
                        except Exception as e:
                            self.logger.debug(f"Image OCR error: {str(e)}")
                            continue
                            
                except Exception as e:
                    self.logger.debug(f"Page image extraction error: {str(e)}")
                    continue
            
            pdf_doc.close()
            
        except Exception as e:
            self.logger.error(f"❌ Image extraction failed: {str(e)}")
        
        return "\n\n".join(image_texts), image_count
    
    def _clean_extracted_text(self, text: str) -> str:
        """Clean extracted text from PDF"""
        # Remove excessive whitespace
        text = re.sub(r'\n{3,}', '\n\n', text)
        text = re.sub(r'[ \t]{2,}', ' ', text)
        
        # Remove page numbers and headers (common patterns)
        text = re.sub(r'\n\d+\s*\n', '\n', text)
        
        # Fix common OCR errors
        text = text.replace('|', 'I')  # Common OCR mistake
        text = text.replace('¢', 'c')
        
        return text.strip()
    
    def _process_pptx(self, content: bytes, file_info: Dict) -> Tuple[str, Dict]:
        """Process PowerPoint files (.pptx)"""
        try:
            # Save to temporary file
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
                temp_file.write(content)
                temp_path = temp_file.name
            
            try:
                prs = Presentation(temp_path)
                slides_text = []
                
                for slide_num, slide in enumerate(prs.slides):
                    slide_text = []
                    
                    # Extract text from shapes
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    
                    if slide_text:
                        slides_text.append(f"--- Slide {slide_num + 1} ---\n" + "\n".join(slide_text))
                
                text_content = "\n\n".join(slides_text)
                metadata = {
                    'slides': len(prs.slides),
                    'extraction_method': 'python-pptx'
                }
                
                return text_content, metadata
                
            finally:
                os.unlink(temp_path)
                
        except Exception as e:
            self.logger.error(f"❌ PPTX processing failed: {str(e)}")
            return "", {}
    
    def _process_docx(self, content: bytes, file_info: Dict) -> Tuple[str, Dict]:
        """Process Word documents (.docx)"""
        try:
            # Save to temporary file
            with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as temp_file:
                temp_file.write(content)
                temp_path = temp_file.name
            
            try:
                doc = Document(temp_path)
                paragraphs = []
                
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        paragraphs.append(paragraph.text.strip())
                
                # Extract text from tables
                for table in doc.tables:
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            paragraphs.append(" | ".join(row_text))
                
                text_content = "\n\n".join(paragraphs)
                metadata = {
                    'paragraphs': len(doc.paragraphs),
                    'tables': len(doc.tables),
                    'extraction_method': 'python-docx'
                }
                
                return text_content, metadata
                
            finally:
                os.unlink(temp_path)
                
        except Exception as e:
            self.logger.error(f"❌ DOCX processing failed: {str(e)}")
            return "", {}
    
    def _process_mp4(self, content: bytes, file_info: Dict) -> Tuple[str, Dict]:
        """Process MP4 video files (extract metadata only for now)"""
        try:
            # For now, we'll just create a placeholder with metadata
            # In future, we could add video transcription capabilities
            filename = file_info['filename']
            size_mb = len(content) / (1024 * 1024)
            
            text_content = f"""
            Video File: {filename}
            File Type: MP4 Video
            Size: {size_mb:.2f} MB
            Grade: {file_info['grade']}
            
            This is a video educational resource. Video content analysis and transcription 
            capabilities can be added in future iterations.
            """
            
            metadata = {
                'file_type': 'video',
                'size_mb': size_mb,
                'extraction_method': 'metadata_only'
            }
            
            return text_content.strip(), metadata
            
        except Exception as e:
            self.logger.error(f"❌ MP4 processing failed: {str(e)}")
            return "", {}
    
    def _process_ppt(self, content: bytes, file_info: Dict) -> Tuple[str, Dict]:
        """Process legacy PowerPoint files (.ppt)"""
        try:
            # For .ppt files, we'll extract what we can or provide metadata
            filename = file_info['filename']
            size_mb = len(content) / (1024 * 1024)
            
            text_content = f"""
            PowerPoint File: {filename}
            File Type: Legacy PowerPoint (.ppt)
            Size: {size_mb:.2f} MB
            Grade: {file_info['grade']}
            
            This is a legacy PowerPoint presentation. For full text extraction,
            consider converting to .pptx format.
            """
            
            metadata = {
                'file_type': 'presentation',
                'size_mb': size_mb,
                'extraction_method': 'metadata_only'
            }
            
            return text_content.strip(), metadata
            
        except Exception as e:
            self.logger.error(f"❌ PPT processing failed: {str(e)}")
            return "", {}
    
    def _process_xlsx(self, content: bytes, file_info: Dict) -> Tuple[str, Dict]:
        """Process Excel files (.xlsx)"""
        try:
            import pandas as pd
            
            # Save to temporary file
            with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as temp_file:
                temp_file.write(content)
                temp_path = temp_file.name
            
            try:
                # Read all sheets
                excel_file = pd.ExcelFile(temp_path)
                sheets_text = []
                
                for sheet_name in excel_file.sheet_names:
                    df = pd.read_excel(temp_path, sheet_name=sheet_name)
                    
                    # Convert to text representation
                    if not df.empty:
                        sheet_text = f"--- Sheet: {sheet_name} ---\n"
                        sheet_text += df.to_string(index=False)
                        sheets_text.append(sheet_text)
                
                text_content = "\n\n".join(sheets_text)
                metadata = {
                    'sheets': len(excel_file.sheet_names),
                    'extraction_method': 'pandas'
                }
                
                return text_content, metadata
                
            finally:
                os.unlink(temp_path)
                
        except ImportError:
            # Fallback if pandas not available
            filename = file_info['filename']
            text_content = f"Excel File: {filename}\nRequires pandas for full text extraction."
            metadata = {'extraction_method': 'metadata_only'}
            return text_content, metadata
        except Exception as e:
            self.logger.error(f"❌ XLSX processing failed: {str(e)}")
            return "", {}
    
    def _process_doc(self, content: bytes, file_info: Dict) -> Tuple[str, Dict]:
        """Process legacy Word documents (.doc)"""
        try:
            # For now, provide metadata only
            # Full .doc processing would require additional libraries like python-docx2txt
            filename = file_info['filename']
            size_mb = len(content) / (1024 * 1024)
            
            text_content = f"""
            Word Document: {filename}
            File Type: Legacy Word Document (.doc)
            Size: {size_mb:.2f} MB
            Grade: {file_info['grade']}
            
            This is a legacy Word document. For full text extraction,
            consider converting to .docx format.
            """
            
            metadata = {
                'file_type': 'document',
                'size_mb': size_mb,
                'extraction_method': 'metadata_only'
            }
            
            return text_content.strip(), metadata
            
        except Exception as e:
            self.logger.error(f"❌ DOC processing failed: {str(e)}")
            return "", {}
    
    def _process_html(self, content: bytes, file_info: Dict) -> Tuple[str, Dict]:
        """Process HTML files"""
        try:
            from bs4 import BeautifulSoup
            
            # Parse HTML
            soup = BeautifulSoup(content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Get text
            text_content = soup.get_text()
            
            # Clean up text
            lines = (line.strip() for line in text_content.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text_content = '\n'.join(chunk for chunk in chunks if chunk)
            
            metadata = {
                'extraction_method': 'beautifulsoup'
            }
            
            return text_content, metadata
            
        except ImportError:
            # Fallback without BeautifulSoup
            text_content = content.decode('utf-8', errors='ignore')
            # Basic HTML tag removal
            text_content = re.sub(r'<[^>]+>', '', text_content)
            metadata = {'extraction_method': 'basic_html'}
            return text_content, metadata
        except Exception as e:
            self.logger.error(f"❌ HTML processing failed: {str(e)}")
            return "", {}
    
    def _create_chunks(self, text: str, file_info: Dict, metadata: Dict, namespace: str) -> List[Dict]:
        """Create chunks from extracted text"""
        if not text or len(text.strip()) < self.min_chunk_length:
            return []
        
        chunks = []
        text = text.strip()
        
        # Split text into chunks
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        
        current_chunk = ""
        current_tokens = 0
        
        for para in paragraphs:
            # Calculate tokens
            if self.tokenizer:
                para_tokens = len(self.tokenizer.encode(para))
            else:
                para_tokens = len(para) // 4
            
            # Check if adding this paragraph would exceed chunk size
            if current_tokens + para_tokens > self.chunk_size and current_chunk:
                # Create chunk
                if len(current_chunk) >= self.min_chunk_length:
                    chunk = self._create_chunk_metadata(
                        current_chunk, file_info, metadata, namespace, len(chunks)
                    )
                    chunks.append(chunk)
                
                # Start new chunk with overlap
                overlap = self._get_overlap(current_chunk)
                current_chunk = overlap + "\n\n" + para if overlap else para
                current_tokens = len(self.tokenizer.encode(current_chunk)) if self.tokenizer else len(current_chunk) // 4
            else:
                current_chunk += "\n\n" + para if current_chunk else para
                current_tokens += para_tokens
        
        # Add final chunk
        if current_chunk and len(current_chunk) >= self.min_chunk_length:
            chunk = self._create_chunk_metadata(
                current_chunk, file_info, metadata, namespace, len(chunks)
            )
            chunks.append(chunk)
        
        self.stats['chunks_created'] += len(chunks)
        return chunks
    
    def _get_overlap(self, text: str) -> str:
        """Get overlap text from the end of current chunk"""
        if len(text) <= self.chunk_overlap:
            return text
        
        overlap_text = text[-self.chunk_overlap:]
        # Try to break at sentence boundary
        sentence_break = overlap_text.find('. ')
        if sentence_break > 0:
            return overlap_text[sentence_break + 2:]
        
        return overlap_text
    
    def _create_chunk_metadata(self, chunk_text: str, file_info: Dict, file_metadata: Dict, namespace: str, chunk_index: int) -> Dict:
        """Create comprehensive Edify-compatible metadata for a chunk"""
        filename = file_info['filename']
        chunk_id = f"{filename}_{chunk_index:03d}"
        preview = chunk_text[:200] if len(chunk_text) > 200 else chunk_text
        
        # Determine content type
        content_type = self._identify_content_type(chunk_text)
        
        # Extract grade and department info
        grade = file_info['grade'].lower()
        department, sub_department = self._get_department_info(grade)
        
        # School types based on namespace
        school_types = []
        if namespace == "kb-esp":
            school_types = ["esp"]
        elif namespace == "kb-psp":
            school_types = ["psp"]
        elif namespace == "kb-msp":
            school_types = ["msp"]
        elif namespace == "kb-ssp":
            school_types = ["ssp"]
        
        # Timestamps
        current_time = time.time()
        current_iso = datetime.now(timezone.utc).isoformat()
        
        # Calculate tokens
        chunk_tokens = len(self.tokenizer.encode(chunk_text)) if self.tokenizer else len(chunk_text) // 4
        
        # Get display title
        display_title = self._get_display_title(filename, grade)
        
        metadata = {
            # IDs (Edify compatibility)
            'id': chunk_id,
            'vector_id': chunk_id,
            'chunk_id': chunk_id,
            'text': chunk_text,
            
            # Chunk information
            'chunk_index': str(chunk_index),
            'section_index': str(0),  # Can be enhanced with section detection
            'chunk_length': str(len(chunk_text)),
            'chunk_tokens': str(chunk_tokens),
            
            # Content fields
            'document_content': chunk_text,
            'preview': preview,
            'content_type': content_type,
            
            # File metadata
            'filename': file_info['blob_name'],
            'original_filename': filename,
            'display_title': display_title,
            'document_type': file_info['file_type'].lstrip('.'),
            'file_type': file_info['file_type'].lstrip('.'),
            'file_pages': str(file_metadata.get('pages', 0)),
            'file_size': str(file_info.get('size', 0) / (1024 * 1024)),
            'file_url': "",
            
            # Processing metadata
            'extraction_method': file_metadata.get('extraction_method', 'unknown'),
            'ocr_used': str(file_metadata.get('ocr_used', False)),
            'images_processed': str(file_metadata.get('images_processed', 0)),
            'repaired': str(file_metadata.get('repaired', False)),
            
            # Department & Grade
            'grade': grade,
            'department': department,
            'sub_department': sub_department,
            'namespace': namespace,
            'school_types': school_types,
            
            # Source tracking
            'source_path': file_info['blob_name'],
            'source_collection': 'enhanced_kb_processor',
            'metadata_source': 'enhanced_kb_processor',
            'has_edify_metadata': True,
            
            # Timestamps
            'created_at': str(current_time),
            'stored_at': str(current_time),
            'enhanced_at': current_iso,
            'last_modified': file_info['last_modified'].isoformat() if file_info.get('last_modified') else current_iso
        }
        
        # Add video-specific metadata if available
        if file_metadata.get('video_url'):
            metadata['video_url'] = file_metadata['video_url']
            metadata['media_type'] = file_metadata.get('media_type', 'video')
            metadata['has_video_content'] = file_metadata.get('has_video_content', True)
            metadata['transcription_available'] = file_metadata.get('transcription_available', False)
            metadata['video_duration'] = file_metadata.get('duration_seconds', 0)
            
        return metadata
    
    def _identify_content_type(self, text: str) -> str:
        """Identify content type from text"""
        text_lower = text.lower()
        
        # Curriculum keywords
        curriculum_keywords = ['lesson', 'chapter', 'exercise', 'question', 'activity', 'assignment']
        if any(kw in text_lower for kw in curriculum_keywords):
            return 'curriculum'
        
        # Conceptual keywords
        conceptual_keywords = ['definition', 'concept', 'theory', 'principle', 'introduction']
        if any(kw in text_lower for kw in conceptual_keywords):
            return 'conceptual'
        
        # Procedural keywords
        procedural_keywords = ['step', 'process', 'procedure', 'method', 'instruction']
        if any(kw in text_lower for kw in procedural_keywords):
            return 'procedural'
        
        # Data visualization keywords
        data_keywords = ['table', 'figure', 'chart', 'graph', 'diagram']
        if any(kw in text_lower for kw in data_keywords):
            return 'data_visualization'
        
        return 'general'
    
    def _get_department_info(self, grade: str) -> Tuple[str, str]:
        """Get department information from grade"""
        grade_lower = grade.lower()
        
        if grade_lower in ['playgroup', 'ik1', 'ik2', 'ik3']:
            return f"{grade.upper()} Curriculum", grade.upper()
        elif 'grade' in grade_lower:
            grade_num = grade_lower.replace('grade', '')
            return f"GRADE{grade_num} Curriculum", f"GRADE{grade_num}"
        elif 'common' in grade_lower:
            return "COMMON Curriculum", "COMMON"
        else:
            return f"{grade.upper()} Curriculum", grade.upper()
    
    def _get_display_title(self, filename: str, grade: str) -> str:
        """Generate display title for file"""
        # Remove extension
        name_without_ext = os.path.splitext(filename)[0]
        
        # Check if it's a UUID-like filename
        if len(name_without_ext) >= 8 and '-' in name_without_ext:
            # Extract first 8 characters
            short_id = name_without_ext[:8]
            return f"{grade.upper()} Document ({short_id})"
        
        # Regular filename
        return name_without_ext.replace('_', ' ').replace('-', ' ').title()
    
    async def _upload_to_pinecone(self, chunks: List[Dict], namespace: str):
        """Upload chunks to Pinecone with embeddings"""
        try:
            if not chunks:
                return
            
            self.logger.info(f"🔄 Generating embeddings for {len(chunks)} chunks...")
            
            # Generate embeddings
            texts = [chunk['text'] for chunk in chunks]
            embeddings = self.embedding_model.encode(texts, convert_to_tensor=False)
            
            # Prepare vectors for Pinecone
            vectors = []
            for chunk, embedding in zip(chunks, embeddings):
                # Create metadata with text field included
                metadata = {}
                for k, v in chunk.items():
                    if k == 'text':
                        # Include full text in metadata (Pinecone supports up to 40KB per metadata field)
                        metadata['text'] = v
                    elif k != 'id':  # Skip 'id' as it's at vector level
                        metadata[k] = v
                
                vector = {
                    'id': chunk['id'],
                    'values': embedding.tolist(),
                    'metadata': metadata
                }
                vectors.append(vector)
            
            # Upload to Pinecone in batches
            batch_size = 100
            for i in range(0, len(vectors), batch_size):
                batch = vectors[i:i + batch_size]
                self.index.upsert(vectors=batch, namespace=namespace)
                
                self.logger.info(f"✅ Uploaded batch {i//batch_size + 1}/{(len(vectors)-1)//batch_size + 1} to {namespace}")
            
            self.logger.info(f"✅ Uploaded {len(chunks)} chunks to namespace: {namespace}")
            return True
            
        except Exception as e:
            self.logger.error(f"❌ Error uploading to Pinecone: {str(e)}")
            return False
    
    def _print_live_statistics(self):
        """Print live processing statistics"""
        elapsed_time = time.time() - self.stats['start_time']
        processed = self.stats['processed_files']
        failed = self.stats['failed_files']
        total = self.stats['total_files']
        
        if total > 0:
            success_rate = (processed / (processed + failed)) * 100 if (processed + failed) > 0 else 0
            completion = self.stats['completion_percentage']
            
            print(f"\n📊 LIVE STATISTICS")
            print(f"═══════════════════════════════════════")
            print(f"📁 Current File: {self.stats['current_file'][:50]}...")
            print(f"⏱️  Elapsed Time: {elapsed_time:.1f}s")
            print(f"📈 Progress: {completion:.1f}% ({processed + failed}/{total})")
            print(f"✅ Success Rate: {success_rate:.1f}%")
            print(f"📄 Chunks Created: {self.stats['chunks_created']:,}")
            print(f"═══════════════════════════════════════")
    
    def _print_final_statistics(self):
        """Print final processing statistics"""
        elapsed_time = time.time() - self.stats['start_time']
        processed = self.stats['processed_files']
        failed = self.stats['failed_files']
        total = self.stats['total_files']
        
        print(f"\n🎉 FINAL PROCESSING STATISTICS")
        print(f"═══════════════════════════════════════════════")
        print(f"⏱️  Total Processing Time: {elapsed_time:.1f}s")
        print(f"📁 Total Files: {total:,}")
        print(f"✅ Successfully Processed: {processed:,}")
        print(f"❌ Failed: {failed:,}")
        print(f"📈 Success Rate: {(processed/total)*100:.1f}%")
        print(f"📄 Total Chunks Created: {self.stats['chunks_created']:,}")
        print(f"⚡ Average Processing Speed: {processed/(elapsed_time/60):.1f} files/minute")
        
        print(f"\n📊 FILES BY TYPE:")
        for file_type, count in self.stats['files_by_type'].items():
            print(f"   {file_type}: {count:,} files")
        
        print(f"\n🎯 FILES BY NAMESPACE:")
        for namespace, count in self.stats['files_by_namespace'].items():
            print(f"   {namespace}: {count:,} files")
        
        print(f"\n⏱️  AVERAGE PROCESSING TIMES BY TYPE:")
        for file_type, times in self.stats['processing_times'].items():
            avg_time = sum(times) / len(times)
            print(f"   {file_type}: {avg_time:.2f}s")
        
        print(f"═══════════════════════════════════════════════")
    
    def process_blob_created(self, container_name: str, blob_name: str, namespace_override: str = "") -> Dict:
        """Process a new/updated blob - full ingestion pipeline (for API)"""
        result = {
            'success': False,
            'blob_name': blob_name,
            'chunks': 0,
            'error': None
        }
        
        start_time = time.time()
        
        try:
            # Use namespace override if provided, otherwise auto-detect
            if namespace_override:
                namespace = namespace_override
                self.logger.info(f"Using provided namespace: {namespace}")
            else:
                namespace = self.get_namespace_from_blob_path(blob_name)
                self.logger.info(f"Auto-detected namespace: {namespace}")
            
            # Extract file info from blob path
            path_parts = blob_name.split('/')
            filename = path_parts[-1]
            file_ext = os.path.splitext(filename)[1].lower()
            
            # Extract grade if available (for KB files)
            grade = path_parts[2] if len(path_parts) >= 3 else 'unknown'
            year = path_parts[1] if len(path_parts) >= 2 else 'unknown'
            
            # Download blob
            container_client = self.blob_service_client.get_container_client(container_name)
            blob_client = container_client.get_blob_client(blob_name)
            file_content = blob_client.download_blob().readall()
            
            # Get blob properties
            blob_properties = blob_client.get_blob_properties()
            
            # Create file_info
            file_info = {
                'filename': filename,
                'blob_name': blob_name,
                'grade': grade,
                'file_type': file_ext,
                'size': len(file_content),
                'last_modified': blob_properties.last_modified
            }
            
            # Get the appropriate handler
            handler = self.file_handlers.get(file_ext)
            if not handler:
                result['error'] = f"No handler for file type: {file_ext}"
                return result
            
            # Process with appropriate handler
            text_content, metadata = handler(file_content, file_info)
            
            if not text_content or len(text_content.strip()) < 50:
                result['error'] = "No text extracted from file"
                return result
            
            # Create chunks
            chunks = self._create_chunks(text_content, file_info, metadata, namespace)
            if not chunks:
                result['error'] = "No chunks created"
                return result
            
            # Delete existing vectors (for updates)
            self.delete_from_pinecone(blob_name, namespace)
            
            # Upload to Pinecone
            asyncio.run(self._upload_to_pinecone(chunks, namespace))
            
            result['success'] = True
            result['chunks'] = len(chunks)
            result['processing_time'] = time.time() - start_time
            result['namespace'] = namespace
            
            self.logger.info(f"✅ Successfully processed {blob_name}: {len(chunks)} chunks in namespace '{namespace}' ({result['processing_time']:.2f}s)")
            
        except Exception as e:
            result['error'] = str(e)
            self.logger.error(f"❌ Processing failed for {blob_name}: {str(e)}")
        
        return result
    
    def delete_from_pinecone(self, blob_name: str, namespace: str = "") -> bool:
        """Delete vectors for a specific file from Pinecone"""
        try:
            # Extract just the filename
            filename_only = os.path.basename(blob_name)
            
            # Auto-detect namespace from blob path if not provided
            if not namespace:
                namespace = self.get_namespace_from_blob_path(blob_name)
            
            self.logger.info(f"Attempting to delete vectors for '{filename_only}' from namespace: '{namespace}'")
            
            # Delete by original_filename (just the filename without path)
            self.index.delete(
                filter={"original_filename": {"$eq": filename_only}},
                namespace=namespace
            )
            
            self.logger.info(f"✅ Deleted vectors for {blob_name} (filename: {filename_only}) from namespace: {namespace}")
            return True
            
        except Exception as e:
            self.logger.error(f"❌ Pinecone delete failed: {str(e)}")
            import traceback
            self.logger.error(traceback.format_exc())
            return False
    
    # ============================================================
    # LOCAL FILE SYSTEM METHODS
    # ============================================================
    
    def load_local_pdf(self, file_path: str) -> Optional[io.BytesIO]:
        """Load PDF from local file system"""
        try:
            if not os.path.exists(file_path):
                self.logger.error(f"File not found: {file_path}")
                return None
            
            with open(file_path, 'rb') as f:
                pdf_stream = io.BytesIO(f.read())
            
            # Verify it's a PDF
            pdf_stream.seek(0)
            header = pdf_stream.read(4)
            if header != b'%PDF':
                self.logger.error(f"Not a valid PDF: {file_path}")
                return None
            
            pdf_stream.seek(0)
            self.logger.info(f"Loaded local PDF: {file_path}")
            return pdf_stream
            
        except Exception as e:
            self.logger.error(f"Failed to load {file_path}: {str(e)}")
            return None
    
    def process_local_file(self, folder_path: str, file_name: str, namespace_override: str = "") -> Dict:
        """Process a local PDF file - full ingestion pipeline"""
        result = {
            'success': False,
            'file_name': file_name,
            'chunks': 0,
            'error': None
        }
        
        start_time = time.time()
        file_path = os.path.join(folder_path, file_name)
        
        try:
            # Use namespace override if provided, otherwise auto-detect
            if namespace_override:
                namespace = namespace_override
                self.logger.info(f"Using provided namespace: {namespace}")
            else:
                namespace = self.get_namespace_from_blob_path(file_path)
                self.logger.info(f"Auto-detected namespace from path: {namespace}")
            
            # Extract file extension
            file_ext = os.path.splitext(file_name)[1].lower()
            
            # 1. Load local file
            if file_ext == '.pdf':
                file_stream = self.load_local_pdf(file_path)
            else:
                # For other formats, load as binary
                with open(file_path, 'rb') as f:
                    file_stream = io.BytesIO(f.read())
            
            if not file_stream:
                result['error'] = f"Failed to load file: {file_path}"
                return result
            
            # Create file_info
            file_info = {
                'filename': file_name,
                'blob_name': file_path,
                'grade': 'unknown',
                'file_type': file_ext,
                'size': os.path.getsize(file_path),
                'last_modified': datetime.fromtimestamp(os.path.getmtime(file_path))
            }
            
            # Get the appropriate handler
            handler = self.file_handlers.get(file_ext)
            if not handler:
                result['error'] = f"No handler for file type: {file_ext}"
                return result
            
            # 2. Extract text with appropriate handler
            text_content, metadata = handler(file_stream.getvalue(), file_info)
            if not text_content or len(text_content.strip()) < 50:
                result['error'] = "No text extracted from file"
                return result
            
            # 3. Create chunks
            chunks = self._create_chunks(text_content, file_info, metadata, namespace)
            if not chunks:
                result['error'] = "No chunks created"
                return result
            
            # 4. Delete existing vectors (for updates)
            self.delete_from_pinecone(file_name, namespace)
            
            # 5. Upload to Pinecone
            asyncio.run(self._upload_to_pinecone(chunks, namespace))
            
            result['success'] = True
            result['chunks'] = len(chunks)
            result['processing_time'] = time.time() - start_time
            result['namespace'] = namespace
            
            self.logger.info(f"✅ Successfully processed local file {file_name}: {len(chunks)} chunks in namespace '{namespace}' ({result['processing_time']:.2f}s)")
            
        except Exception as e:
            result['error'] = str(e)
            self.logger.error(f"❌ Processing failed for local file {file_name}: {str(e)}")
            import traceback
            self.logger.error(traceback.format_exc())
        
        return result
    
    def process_local_file_deleted(self, file_name: str, namespace: str = "") -> Dict:
        """Remove a local file from Pinecone"""
        result = {
            'success': False,
            'file_name': file_name,
            'error': None
        }
        
        try:
            success = self.delete_from_pinecone(file_name, namespace)
            result['success'] = success
            if not success:
                result['error'] = "Failed to delete from Pinecone"
            else:
                self.logger.info(f"✅ Successfully removed local file {file_name} from Pinecone")
                
        except Exception as e:
            result['error'] = str(e)
            self.logger.error(f"❌ Delete failed for local file {file_name}: {str(e)}")
        
        return result

# Async wrapper for the main processing
async def main():
    """Main processing function"""
    import os
    from dotenv import load_dotenv
    
    # Load environment variables
    load_dotenv()
    
    # Configuration
    config = {
        'chunk_size': int(os.getenv('CHUNK_SIZE', 600)),
        'chunk_overlap': int(os.getenv('CHUNK_OVERLAP', 120)),
        'min_chunk_length': int(os.getenv('MIN_CHUNK_LENGTH', 100)),
        'max_chunk_length': int(os.getenv('MAX_CHUNK_LENGTH', 1200)),
        'batch_size': int(os.getenv('BATCH_SIZE', 15)),
        'max_workers': int(os.getenv('MAX_WORKERS', 4)),
        'enable_ocr': True,
        'enable_repair': True
    }
    
    # Initialize processor
    processor = EnhancedKBProcessor(config)
    
    # Process KB directory (12/ only as specified)
    await processor.process_kb_directory(target_year='12')

if __name__ == "__main__":
    import asyncio
    asyncio.run(main())